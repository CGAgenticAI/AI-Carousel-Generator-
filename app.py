### app.py — Streamlit front‑end for AI‑generated LinkedIn carousels
"""
Minimal Streamlit application that:
1. Lets the user enter a *topic* and a **Groq API key**.
2. Calls Groq via Llama‑Index to create a JSON description of a carousel (cover + Q&A slides).
3. Renders the JSON into a square PowerPoint deck with black background & white text.
4. Presents the deck for download.

Deployment‑ready for **Streamlit Community Cloud**:
- Put this single file in the repo root (e.g. `app.py`).
- Add a `requirements.txt` with:
    streamlit
    python‑pptx
    python‑dotenv
    llama‑index
    groq
- Push to GitHub → Create new app on share.streamlit.io → Add a secret called
  `GROQ_API_KEY` if you prefer not to type it each time.
"""

from __future__ import annotations

import json
import os
from io import BytesIO
from pathlib import Path
from typing import List, Literal

import streamlit as st
from dotenv import load_dotenv
from llama_index.llms.groq import Groq
from pydantic import BaseModel, Field
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ────────────────────────────────────────────────────────────────────────────────
# Pydantic models (same as notebook version)
# ────────────────────────────────────────────────────────────────────────────────
class CoverSlide(BaseModel):
    slide_number: int = Field(description="Position of this slide in the carousel (must be 1)")
    type: Literal["cover"] = Field(description="Always 'cover' for the first slide")
    title: str
    subtitle: str


class QASlide(BaseModel):
    slide_number: int
    type: Literal["qa"] = Field(description="'qa' for every question‑answer slide")
    question: str
    answer: str


class Carousel(BaseModel):
    cover: CoverSlide
    qa_slides: List[QASlide]


# ────────────────────────────────────────────────────────────────────────────────
# Helper: build the analysis prompt
# ────────────────────────────────────────────────────────────────────────────────
PROMPT_TMPL = """You are an expert content strategist.\n\nTopic: {topic}\n\nTask:\n1. Write an inspiring one‑liner subtitle for the cover slide.\n2. Then identify between 5 and 8 key questions someone would ask about this topic.\n3. For each question, provide a concise (2‑3 line) answer.\n\nFormat your reply exactly as:\n\nTopic: {topic}\nSubtitle: <your one‑liner>\n\nQuestions & Answers:\n1. Question: <question 1>\n   Answer: <answer 1>\n\n2. Question: <question 2>\n   Answer: <answer 2>\n…"""


# ────────────────────────────────────────────────────────────────────────────────
# PPTX helpers
# ────────────────────────────────────────────────────────────────────────────────

def _black_bg(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0, 0, 0)


def _style_para(para, font_size: int, bold: bool = False):
    para.font.name = "Calibri"
    para.font.size = Pt(font_size)
    para.font.bold = bold
    para.font.color.rgb = RGBColor(255, 255, 255)
    para.alignment = PP_ALIGN.LEFT


def build_pptx(carousel: Carousel) -> BytesIO:
    """Render a Carousel → PPTX and return as BytesIO."""
    prs = Presentation()
    prs.slide_width = Inches(11.25)
    prs.slide_height = Inches(11.25)

    # Cover slide
    cover = carousel.cover
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _black_bg(slide)

    title_tf = slide.shapes.title.text_frame
    title_tf.clear(); title_tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    title_tf.margin_left = title_tf.margin_right = title_tf.margin_top = title_tf.margin_bottom = 0
    _style_para(title_tf.paragraphs[0], 54, bold=True)
    title_tf.paragraphs[0].text = cover.title

    subt_tf = slide.placeholders[1].text_frame
    subt_tf.clear(); subt_tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    subt_tf.margin_left = subt_tf.margin_right = subt_tf.margin_top = subt_tf.margin_bottom = 0
    _style_para(subt_tf.paragraphs[0], 34)
    subt_tf.paragraphs[0].text = cover.subtitle

    # QA slides
    for qa in carousel.qa_slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        _black_bg(slide)

        q_tf = slide.shapes.title.text_frame
        q_tf.clear(); q_tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        q_tf.margin_left = q_tf.margin_right = q_tf.margin_top = q_tf.margin_bottom = 0
        _style_para(q_tf.paragraphs[0], 40, bold=True)
        q_tf.paragraphs[0].text = qa.question

        a_tf = slide.placeholders[1].text_frame
        a_tf.clear(); a_tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        a_tf.margin_left = a_tf.margin_right = a_tf.margin_top = a_tf.margin_bottom = 0
        _style_para(a_tf.paragraphs[0], 30)
        a_tf.paragraphs[0].text = qa.answer

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ────────────────────────────────────────────────────────────────────────────────
# Streamlit UI
# ────────────────────────────────────────────────────────────────────────────────
st.set_page_config(page_title="AI Carousel Builder", layout="centered")
st.title("📊 LinkedIn Carousel Generator")

with st.form("input_form"):
    topic = st.text_input("🎯 Topic", placeholder="e.g. Game Theory", help="What do you want the carousel to be about?")

    # Prefer Streamlit secrets but allow manual entry for local tests
    api_default = os.getenv("GROQ_API_KEY", "") or st.secrets.get("GROQ_API_KEY", "")
    api_key = st.text_input("🔑 Groq API key", value=api_default, type="password")

    submit = st.form_submit_button("Generate Carousel → PPTX")

if submit:
    if not topic:
        st.error("Please enter a topic.")
        st.stop()
    if not api_key:
        st.error("Please enter your Groq API key.")
        st.stop()

    load_dotenv(override=False)  # allow .env fallback for local dev
    os.environ["GROQ_API_KEY"] = api_key

    # LLM setup
    with st.spinner("Contacting Groq, generating content…"):
        llm = Groq(model="llama-3.3-70b-versatile", api_key=api_key)
        sllm = llm.as_structured_llm(Carousel)
        prompt = PROMPT_TMPL.format(topic=topic)
        raw_response = llm.complete(prompt).text
        carousel: Carousel = sllm.complete(raw_response)

    # Build PPTX
    with st.spinner("Rendering PPTX…"):
        ppt_buf = build_pptx(carousel)
        file_name = f"LinkedIn_Carousel_{topic.replace(' ', '_')}.pptx"

    st.success("Done! Download your deck below.")
    st.download_button(
        label="⬇️ Download PPTX",
        data=ppt_buf,
        file_name=file_name,
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

st.caption("Made with ❤️ & Groq · © 2025")
